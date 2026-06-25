"""pptx_scanner.py — PPT 도형 인벤토리 추출 + 정답지 셀값 읽기.

저장 정보 (모든 도형):
  shape_id  — XML id 속성 (shape_num보다 안정적)
  num       — 이름 끝 숫자 (기존 하위 호환)
  z_order   — 슬라이드 내 렌더링 순서 (0 = 맨 아래)
  left/top/width/height — EMU 단위 위치·크기

표(table):
  grid       — 2D 텍스트 (read_values=True, 하위 호환)
  grid_meta  — 2D 셀별 {text, bold, font_size_pt, fill_type, fill_rgb}

그룹(group):
  children   — 자식 도형 재귀 처리

차트(chart):
  chart_type, categories, series, series_count
"""
from __future__ import annotations

import hashlib
import json
import os
import re
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

_NUM_RE = re.compile(r"(\d+)\s*$")

# 템플릿의 0-값 placeholder 패턴 (외부 모듈에서도 사용 가능)
TEMPLATE_PLACEHOLDER_RE = re.compile(
    r'[+-]?0\.0%p'    # ±0.0%p (MoM)
    r'|0\.0%(?!p)'    # 0.0%   (pct)
    r'|0\.0[KB]'      # 0.0K / 0.0B
    r'|x0\.0'         # x0.0   (ratio)
    r'|x1\.0'         # x1.0   (ratio neutral)
)


def _shape_num(shape) -> int | None:
    m = _NUM_RE.search(shape.name.strip())
    return int(m.group(1)) if m else None


def _emu(v) -> int | None:
    try:
        return int(v)
    except (TypeError, ValueError):
        return None


# ── 표 셀 메타 ────────────────────────────────────────────────────────────────

def _cell_meta(cell) -> dict:
    """셀 하나의 텍스트 + 기본 서식 정보."""
    text = cell.text.strip()
    paras = cell.text_frame.paragraphs
    runs = [r for p in paras for r in p.runs]

    meta: dict[str, Any] = {"text": text}
    if runs:
        run = runs[0]
        meta["bold"] = run.font.bold                                     # True/False/None(상속)
        fs = run.font.size
        meta["font_size_pt"] = round(fs / 12700, 1) if fs else None      # EMU → pt
    else:
        meta["bold"] = None
        meta["font_size_pt"] = None

    try:
        fill = cell.fill
        if fill.type is not None:
            fname = fill.type.name if hasattr(fill.type, "name") else str(fill.type)
            meta["fill_type"] = fname
            if fname == "SOLID":
                try:
                    meta["fill_rgb"] = str(fill.fore_color.rgb)
                except Exception:
                    meta["fill_rgb"] = None
    except Exception:
        pass

    return meta


def _table_to_grid(table) -> list[list[str]]:
    """2D 텍스트 그리드 (하위 호환)."""
    return [
        [table.cell(r, c).text.strip() for c in range(len(table.columns))]
        for r in range(len(table.rows))
    ]


def _table_to_grid_meta(table) -> list[list[dict]]:
    """2D 셀 메타 그리드 (서식 포함)."""
    return [
        [_cell_meta(table.cell(r, c)) for c in range(len(table.columns))]
        for r in range(len(table.rows))
    ]


# ── 차트 ─────────────────────────────────────────────────────────────────────

def _find_chart_shape(slide, shape_num: int | None, shape_id: int | None = None):
    """차트 도형 탐색 (shape_id → shape_num 순)."""
    if shape_id is not None:
        for sh in slide.shapes:
            if sh.shape_id == shape_id and hasattr(sh, "has_chart") and sh.has_chart:
                return sh
    if shape_num is not None:
        for sh in slide.shapes:
            if _shape_num(sh) == shape_num and hasattr(sh, "has_chart") and sh.has_chart:
                return sh
    return None


def write_chart(
    prs,
    slide_idx: int,
    shape_num: int,
    categories: list,
    series: list[tuple],
    *,
    shape_id: int | None = None,
) -> bool:
    """차트 데이터 교체. series = [(name, [v0, v1, ...]), ...].

    None 값은 기존 카테고리 값을 유지하기 어려우므로 0.0으로 대체(빈 막대).
    python-pptx 권장 방식인 chart.replace_data(CategoryChartData) 사용.
    """
    try:
        from pptx.chart.data import CategoryChartData
        sh = _find_chart_shape(prs.slides[slide_idx], shape_num, shape_id)
        if not sh:
            return False
        cd = CategoryChartData()
        cd.categories = [str(c) for c in categories]
        for name, vals in series:
            clean = [float(v) if isinstance(v, (int, float)) else 0.0 for v in vals]
            cd.add_series(str(name), clean)
        sh.chart.replace_data(cd)
        return True
    except Exception:
        return False


def _chart_info(shape) -> dict:
    """차트 타입·카테고리·시리즈 추출."""
    try:
        chart = shape.chart
        try:
            cats = [str(c) for c in list(chart.plots[0].categories)]
        except Exception:
            cats = []
        series = []
        for s in chart.series:
            vals = [round(v, 4) if v is not None else None for v in s.values]
            series.append({"name": s.name, "values": vals})
        raw_type = str(chart.chart_type)
        m = re.search(r"ChartType\.(\w+)", raw_type)
        return {
            "chart_type": m.group(1) if m else raw_type,
            "categories": cats,
            "series": series,
            "series_count": len(series),
        }
    except Exception as e:
        return {"error": str(e)}


# ── 도형 스캔 (재귀) ──────────────────────────────────────────────────────────

def _scan_shape(sh, z_order: int, read_values: bool) -> dict:
    """도형 하나의 전체 정보 (그룹은 children 재귀)."""
    info: dict[str, Any] = {
        "name": sh.name,
        "shape_id": sh.shape_id,
        "num": _shape_num(sh),
        "z_order": z_order,
        "left": _emu(sh.left),
        "top": _emu(sh.top),
        "width": _emu(sh.width),
        "height": _emu(sh.height),
        "type": "other",
    }

    # 그룹 — 자식 재귀
    if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
        info["type"] = "group"
        children = []
        try:
            for ci, child in enumerate(sh.shapes):
                children.append(_scan_shape(child, ci, read_values))
        except Exception:
            pass
        info["children"] = children

    # 표
    elif sh.shape_type == MSO_SHAPE_TYPE.TABLE or (hasattr(sh, "has_table") and sh.has_table):
        info["type"] = "table"
        tbl = sh.table
        info["rows"] = len(tbl.rows)
        info["cols"] = len(tbl.columns)
        if read_values:
            info["grid"] = _table_to_grid(tbl)           # 텍스트만 (하위 호환)
            info["grid_meta"] = _table_to_grid_meta(tbl)  # 서식 포함

    # 차트
    elif hasattr(sh, "has_chart") and sh.has_chart:
        info["type"] = "chart"
        if read_values:
            info["chart"] = _chart_info(sh)

    # 텍스트 도형
    elif hasattr(sh, "has_text_frame") and sh.has_text_frame:
        info["type"] = "text"
        if read_values:
            info["text"] = sh.text_frame.text.strip()[:200]

    return info


# ── Smart Block 그루핑 ────────────────────────────────────────────────────────
# 공간 근접성 임계값 (EMU 단위, 1cm ≈ 360,000)
_BLOCK_VERT_THRESHOLD  = 700_000   # ~1.9cm — label과 table 최대 수직 거리
_BLOCK_HORIZ_THRESHOLD = 6_000_000 # ~16.7cm — 동일 수평 밴드 판정 범위


def group_blocks_from_scan(shapes: list[dict]) -> list[dict]:
    """scan_slide 출력 shapes에서 semantic block 목록을 생성한다.

    block_type:
      KPI_CARD     — 소형 표(≤5행) + 인접 텍스트 레이블 쌍
      TABLE_SECTION — 대형 표(>5행) 단독 (KPI 격자 등)
      CHART        — 차트 도형
      LABEL        — 미분류 텍스트 (섹션 헤더 등)
    """
    texts  = [s for s in shapes if s.get("type") == "text" and s.get("text")]
    tables = [s for s in shapes if s.get("type") == "table"]
    charts = [s for s in shapes if s.get("type") == "chart"]

    blocks: list[dict] = []
    used_label_nums: set = set()

    # ── 1. CHART 블록 ──────────────────────────────────────────────────────
    for ch in charts:
        blocks.append({
            "block_type": "CHART",
            "shape_num":  ch.get("num"),
            "shape_id":   ch.get("shape_id"),
        })

    # ── 2. KPI_CARD: 소형 표 + 인접 레이블 ───────────────────────────────
    small = sorted(
        [t for t in tables if (t.get("rows") or 0) <= 5],
        key=lambda s: (s.get("top") or 0),
    )
    large = [t for t in tables if (t.get("rows") or 0) > 5]

    for tbl in small:
        t_top  = tbl.get("top")  or 0
        t_left = tbl.get("left") or 0

        best: dict | None = None
        best_score = float("inf")
        for txt in texts:
            if txt.get("num") in used_label_nums:
                continue
            if is_kpi_placeholder(txt.get("text", "")):
                continue          # KPI 수치 자체는 레이블이 아님
            vert  = abs((txt.get("top")  or 0) - t_top)
            horiz = abs((txt.get("left") or 0) - t_left)
            if vert < _BLOCK_VERT_THRESHOLD and horiz < _BLOCK_HORIZ_THRESHOLD:
                score = vert + horiz * 0.15   # 수직 거리에 더 가중치
                if score < best_score:
                    best, best_score = txt, score

        block: dict = {
            "block_type": "KPI_CARD",
            "table_num":  tbl.get("num"),
            "table_id":   tbl.get("shape_id"),
            "label":      None,
            "label_num":  None,
        }
        if best:
            block["label"]     = best.get("text", "").strip()
            block["label_num"] = best.get("num")
            used_label_nums.add(best.get("num"))
        blocks.append(block)

    # ── 3. TABLE_SECTION: 대형 표 ─────────────────────────────────────────
    for tbl in sorted(large, key=lambda s: (s.get("top") or 0)):
        blocks.append({
            "block_type": "TABLE_SECTION",
            "table_num":  tbl.get("num"),
            "table_id":   tbl.get("shape_id"),
            "rows":       tbl.get("rows"),
            "cols":       tbl.get("cols"),
        })

    # ── 4. 미분류 텍스트 → LABEL ──────────────────────────────────────────
    for txt in sorted(texts, key=lambda s: (s.get("top") or 0)):
        if txt.get("num") in used_label_nums:
            continue
        if is_kpi_placeholder(txt.get("text", "")):
            continue
        blocks.append({
            "block_type": "LABEL",
            "shape_num":  txt.get("num"),
            "text":       txt.get("text", "").strip(),
        })

    return blocks


def scan_slide(slide, slide_idx: int, read_values: bool = False) -> dict:
    """슬라이드 하나의 도형 인벤토리를 반환. read_values=True이면 blocks도 포함."""
    shapes = [_scan_shape(sh, z, read_values) for z, sh in enumerate(slide.shapes)]
    result: dict = {"slide_idx": slide_idx, "shapes": shapes}
    if read_values:
        result["blocks"] = group_blocks_from_scan(shapes)
    return result


def scan_pptx(path: str, read_values: bool = False) -> list[dict]:
    """PPTX 전체를 스캔해 슬라이드별 도형 인벤토리를 반환."""
    prs = Presentation(path)
    return [scan_slide(slide, idx, read_values) for idx, slide in enumerate(prs.slides)]


# ── JSON 중간문서 캐시 (읽기를 JSON으로) ─────────────────────────────────────
_SCANS_DIR = os.path.join(
    os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__)))),
    "cache", "scans")


def _file_sig(path: str) -> str:
    """파일 내용(첫 1MB) + mtime 시그니처."""
    h = hashlib.md5()
    with open(path, "rb") as f:
        h.update(f.read(1024 * 1024))
    h.update(str(os.path.getmtime(path)).encode())
    return h.hexdigest()[:12]


def scan_pptx_cached(path: str, read_values: bool = False) -> list[dict]:
    """PPTX를 1회만 스캔해 JSON으로 캐시하고, 이후엔 JSON을 읽는다.

    같은 파일(내용+mtime 동일)을 여러 단계가 반복 스캔하는 비용을 제거한다.
    XML(OOXML) 파싱은 최초 1회, 그 다음은 순수 JSON 읽기.
    """
    os.makedirs(_SCANS_DIR, exist_ok=True)
    sig = _file_sig(path)
    cache_path = os.path.join(_SCANS_DIR, f"scan_{sig}_{'v' if read_values else 's'}.json")
    if os.path.exists(cache_path):
        try:
            with open(cache_path, encoding="utf-8") as f:
                return json.load(f)
        except Exception:
            pass
    data = scan_pptx(path, read_values)
    try:
        with open(cache_path, "w", encoding="utf-8") as f:
            json.dump(data, f, ensure_ascii=False)
    except Exception:
        pass
    return data


# ── 도형 탐색 헬퍼 ────────────────────────────────────────────────────────────

def is_kpi_placeholder(text: str) -> bool:
    """텍스트 도형 값이 KPI 수치인지 확인 (+0.0%p, 0.0%, 0.0K, x0.0 등)."""
    t = text.strip()
    if not t:
        return False
    return bool(re.search(r'[+\-]?\d+\.?\d*(%p|%|K)|^x\d+\.?\d*$', t))


def _find_table_shape(slide, shape_num: int | None, shape_id: int | None = None):
    """표 도형을 shape_id → shape_num 순으로 탐색.

    shape_id가 주어지면 정확 매칭을 우선 시도.
    템플릿 편집으로 이름 끝 숫자가 바뀌어도 shape_id 기반으로 안정 동작.
    """
    if shape_id is not None:
        for sh in slide.shapes:
            if sh.shape_id == shape_id and hasattr(sh, "has_table") and sh.has_table:
                return sh
    if shape_num is not None:
        for sh in slide.shapes:
            if _shape_num(sh) == shape_num and hasattr(sh, "has_table") and sh.has_table:
                return sh
    return None


def _find_text_shape(slide, shape_num: int | None, shape_id: int | None = None):
    """표가 아닌 텍스트 도형 탐색 (shape_id → shape_num 순)."""
    if shape_id is not None:
        for sh in slide.shapes:
            if (sh.shape_id == shape_id
                    and hasattr(sh, "has_text_frame") and sh.has_text_frame
                    and not (hasattr(sh, "has_table") and sh.has_table)):
                return sh
    if shape_num is not None:
        for sh in slide.shapes:
            if (_shape_num(sh) == shape_num
                    and hasattr(sh, "has_text_frame") and sh.has_text_frame
                    and not (hasattr(sh, "has_table") and sh.has_table)):
                return sh
    return None


# ── 공개 읽기/쓰기 API ────────────────────────────────────────────────────────

def read_cell(
    prs: Presentation,
    slide_idx: int,
    shape_num: int,
    row: int,
    col: int,
    *,
    shape_id: int | None = None,
) -> str | None:
    """특정 셀 텍스트 읽기.

    shape_id 를 넘기면 shape_num 이름 의존을 우회할 수 있다.
    """
    try:
        sh = _find_table_shape(prs.slides[slide_idx], shape_num, shape_id)
        if sh:
            return sh.table.cell(row, col).text.strip()
    except Exception:
        pass
    return None


_BR_TAG = "{http://schemas.openxmlformats.org/drawingml/2006/main}br"
_R_TAG  = "{http://schemas.openxmlformats.org/drawingml/2006/main}r"
_T_TAG  = "{http://schemas.openxmlformats.org/drawingml/2006/main}t"


def _write_para(para, value: str) -> None:
    """단락에 value를 씀.

    <a:br> 가 있는 단락은 br 이후 run만 교체해 'vs. Global' 접두사를 보존한다.
    없으면 첫 번째 run을 교체한다.
    """
    children = list(para._p)
    br_indices = [i for i, e in enumerate(children) if e.tag == _BR_TAG]

    if br_indices:
        last_br_idx = br_indices[-1]
        post_runs = [e for e in children[last_br_idx + 1:] if e.tag == _R_TAG]
        if post_runs:
            for r in post_runs:
                t = r.find(_T_TAG)
                if t is not None:
                    t.text = ""
            t = post_runs[0].find(_T_TAG)
            if t is not None:
                t.text = value
            return

    # br 없는 일반 단락: 모든 run 비우고 첫 run에 값 쓰기
    for run in para.runs:
        run.text = ""
    if para.runs:
        para.runs[0].text = value
    else:
        para.add_run().text = value


def write_cell(
    prs: Presentation,
    slide_idx: int,
    shape_num: int,
    row: int,
    col: int,
    value: str,
    *,
    shape_id: int | None = None,
) -> bool:
    """특정 셀에 텍스트를 씀 (기존 run 서식 보존).

    <a:br> 가 있는 셀(예: 'vs. Global\\n+0.0%p')은 br 뒤 부분만 교체한다.
    shape_id 를 넘기면 shape_num 이름 의존을 우회할 수 있다.
    """
    try:
        sh = _find_table_shape(prs.slides[slide_idx], shape_num, shape_id)
        if sh:
            cell = sh.table.cell(row, col)
            if cell.text_frame.paragraphs:
                _write_para(cell.text_frame.paragraphs[0], value)
            return True
    except Exception:
        pass
    return False


def read_text_shape(
    prs: Presentation,
    slide_idx: int,
    shape_num: int,
    *,
    shape_id: int | None = None,
) -> str | None:
    """텍스트 도형 값 읽기 (표 제외). shape_id 우선."""
    try:
        sh = _find_text_shape(prs.slides[slide_idx], shape_num, shape_id)
        if sh:
            return sh.text_frame.text.strip()
    except Exception:
        pass
    return None


def _write_placeholder_in_tf(tf, value: str) -> bool:
    """텍스트 프레임에서 placeholder를 찾아 정밀 교체 (Placeholder Preservation).

    우선순위:
    1. TEMPLATE_PLACEHOLDER_RE 일치 run → run 내에서 regex sub (인라인 접두 보존, 최우선)
    2. <a:br> 포함 단락 → _write_para (br 뒤 value 교체) — placeholder가 split run인 경우 폴백
    3. 매칭 실패 → False 반환 (호출자가 fallback 처리)
    """
    # 1. placeholder 패턴 일치 run 탐색 (예: "MoM ▲ 0.0%p" 단일 run, "vs. Co.A +0.0%p")
    for para in tf.paragraphs:
        for run in para.runs:
            if TEMPLATE_PLACEHOLDER_RE.search(run.text):
                # 인라인 접두("MoM ▲ ", "vs. Co.A ")를 보존하고 숫자 부분만 교체
                run.text = TEMPLATE_PLACEHOLDER_RE.sub(value, run.text, count=1)
                return True

    # 2. <a:br> 있는 단락 — placeholder가 split run이라 1에서 못 찾은 경우
    for para in tf.paragraphs:
        if any(e.tag == _BR_TAG for e in para._p):
            _write_para(para, value)
            return True

    return False


def write_text_shape(
    prs: Presentation,
    slide_idx: int,
    shape_num: int,
    value: str,
    *,
    shape_id: int | None = None,
) -> bool:
    """텍스트 도형에 값을 씀 (Placeholder Preservation 적용). shape_id 우선.

    - placeholder 패턴이 있으면 해당 run만 교체해 접두 텍스트 보존
    - 없으면(이미 채워진 도형 재쓰기) 마지막 비어있지 않은 단락에 쓰기
    """
    try:
        sh = _find_text_shape(prs.slides[slide_idx], shape_num, shape_id)
        if sh:
            tf = sh.text_frame
            if not tf.paragraphs:
                return True
            # placeholder 기반 정밀 교체 시도
            if _write_placeholder_in_tf(tf, value):
                return True
            # fallback: 마지막 비어있지 않은 단락 (또는 첫 번째 단락)
            paras_with_text = [p for p in tf.paragraphs if p.text.strip()]
            target = paras_with_text[-1] if paras_with_text else tf.paragraphs[0]
            _write_para(target, value)
            return True
    except Exception:
        pass
    return False


_INSIGHT_PH_RE = re.compile(r'\{\{[^}]+\}\}')


def write_insight_shape(
    prs: Presentation,
    slide_idx: int,
    shape_num: int | None,
    sentence: str,
    *,
    shape_id: int | None = None,
) -> bool:
    """인사이트 문장을 텍스트 박스에 기록.

    동작:
    - {{...}} 단락이 있으면 첫 번째를 sentence로 교체, 나머지 지움 (placeholder 교체)
    - {{...}} 없고 텍스트 박스가 비어있으면 첫 단락에 sentence 기록 (빈 박스 채우기)
    - 텍스트가 이미 있고 {{...}} 도 없으면 False 반환 (덮어쓰기 안 함)
    """
    try:
        sh = _find_text_shape(prs.slides[slide_idx], shape_num, shape_id)
        if not sh:
            return False
        tf = sh.text_frame
        ph_paras = [p for p in tf.paragraphs if _INSIGHT_PH_RE.search(p.text)]
        if ph_paras:
            _write_para(ph_paras[0], sentence)
            for p in ph_paras[1:]:
                _write_para(p, "")
            return True
        # 빈 텍스트 박스: 첫 단락에 직접 기록
        if not tf.text.strip():
            _write_para(tf.paragraphs[0], sentence)
            return True
        return False
    except Exception:
        return False
