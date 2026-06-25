"""Filler Skill 4 — Fill Completeness & Integrity Guard 구현.

check_fill_completeness(): filled_count vs mapping_count 검증
scan_unfilled()          : PPT 저장 전 전체 도형에서 잔여 placeholder 탐색
"""
from __future__ import annotations

import re

from pptx import Presentation

# ── Placeholder 패턴 (템플릿의 0-값 패턴) ────────────────────────────────────
_UNFILLED_RE = re.compile(
    r'[+-]?0\.0%p'    # ±0.0%p (MoM)
    r'|0\.0%(?!p)'    # 0.0%   (pct)
    r'|0\.0[KB]'      # 0.0K / 0.0B
    r'|x0\.0'         # x0.0   (ratio)
    r'|x1\.0'         # x1.0   (ratio neutral)
)


def check_fill_completeness(
    mapping_count: int,
    filled_count: int,
) -> tuple[bool, list[str]]:
    """filled_count == mapping_count 인지 검증.

    반환: (passed, 경고 메시지 목록)
    """
    if filled_count == mapping_count:
        return True, []

    short = mapping_count - filled_count
    msg = (
        f"[Skill 4] FAIL Fill Completeness: "
        f"filled={filled_count} / required={mapping_count} (미채우기 {short}개)"
    )
    return False, [msg]


def scan_unfilled(prs: Presentation) -> list[dict]:
    """PPT 저장 전 전체 도형(표·텍스트·차트)에서 미채움 탐색.

    - 표/텍스트: 잔여 placeholder(0.0% 등)
    - 차트: 모든 계열 값이 0/None (미채움)
    반환: [{"slide_idx", "shape", "type", "text", ...}]
    """
    found: list[dict] = []

    for slide_idx, slide in enumerate(prs.slides):
        for sh in slide.shapes:
            _scan_shape(sh, slide_idx, found)

    return found


def _scan_shape(sh, slide_idx: int, found: list[dict]) -> None:
    """단일 도형을 재귀 탐색 (그룹 포함)."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    # 그룹: 자식 재귀
    if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
        try:
            for child in sh.shapes:
                _scan_shape(child, slide_idx, found)
        except Exception:
            pass
        return

    # 테이블
    if hasattr(sh, "has_table") and sh.has_table:
        try:
            tbl = sh.table
            for r in range(len(tbl.rows)):
                for c in range(len(tbl.columns)):
                    text = tbl.cell(r, c).text.strip()
                    if _UNFILLED_RE.search(text):
                        found.append({
                            "slide_idx": slide_idx,
                            "shape": sh.name,
                            "type": "table",
                            "row": r,
                            "col": c,
                            "text": text,
                        })
        except Exception:
            pass
        return

    # 차트: 모든 계열 값이 0/None이면 미채움으로 간주
    if hasattr(sh, "has_chart") and sh.has_chart:
        try:
            vals = [v for s in sh.chart.series for v in s.values]
            if vals and all(v in (0, None) for v in vals):
                found.append({
                    "slide_idx": slide_idx,
                    "shape": sh.name,
                    "type": "chart",
                    "text": "chart(모든 계열값 0/None)",
                })
        except Exception:
            pass
        return

    # 텍스트 도형
    if hasattr(sh, "has_text_frame") and sh.has_text_frame:
        try:
            text = sh.text_frame.text.strip()
            if _UNFILLED_RE.search(text):
                found.append({
                    "slide_idx": slide_idx,
                    "shape": sh.name,
                    "type": "text",
                    "text": text,
                })
        except Exception:
            pass
