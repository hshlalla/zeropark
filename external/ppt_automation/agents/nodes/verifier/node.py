"""④ Verifier 에이전트 — SlideMapping.targets만 검증.

정답지 전체 표 스캔 대신 SlideMapping에 등록된 셀만 대조한다.
문장형 인사이트 표 등 KPI가 아닌 셀은 검증 대상에서 자동 제외.
"""
from __future__ import annotations

import os
import re

from langchain_anthropic import ChatAnthropic
from langchain_core.messages import AIMessage, HumanMessage, SystemMessage
from pptx import Presentation

from agents.state import AgentState
from agents.models import (
    VerificationResult, VerificationIssue,
    SlideMapping, ValidationIssue, FailureReport,
)
from core.predefined.pptx_scanner import read_cell, read_text_shape
from .kpi_conservation import check_conservation, print_conservation_report
from agents.utils import load_contract, get_anthropic_api_key
from agents.nodes.manager import HARD_ITERATION_LIMIT
from domain.config import TOLERANCE_PCT, TOLERANCE_MOM, TOLERANCE_OTHER

_CONTRACT = load_contract("verifier")   # agents/contracts/verifier.md

_NUMERIC_RE = re.compile(r"[-+]?[\d.]+")
# 반복 상한은 Manager의 HARD_ITERATION_LIMIT가 단일 출처. Verifier는 이를 참조만 한다.
MAX_RETRIES = HARD_ITERATION_LIMIT


def _strip_prefix(s: str) -> str:
    """'vs. Global\\x0b+3.6%p' → '+3.6%p' (소프트 개행 앞 레이블 제거)."""
    return s.split("\x0b")[-1] if "\x0b" in s else s


def _to_float(s: str) -> float | None:
    s = _strip_prefix(s)
    s = s.replace("K", "").replace("%p", "").replace("%", "").strip()
    m = _NUMERIC_RE.search(s)
    try:
        return float(m.group()) if m else None
    except ValueError:
        return None


def _is_close(a: str, b: str, tol: float = 0.15) -> bool:
    # 소프트 개행 접두사 제거 후 비교 (vs. Global\x0b+3.6%p → +3.6%p)
    a_cmp = _strip_prefix(a)
    b_cmp = _strip_prefix(b)
    if a_cmp.strip() == b_cmp.strip():
        return True
    if a_cmp.strip() in ("-", "") and b_cmp.strip() in ("-", ""):
        return True
    fa, fb = _to_float(a), _to_float(b)
    if fa is not None and fb is not None:
        a_is_pct = "%" in a_cmp and "%p" not in a_cmp
        b_is_pct = "%" in b_cmp and "%p" not in b_cmp
        a_is_mom = "%p" in a_cmp
        b_is_mom = "%p" in b_cmp
        if a_is_pct or b_is_pct:
            effective_tol = TOLERANCE_PCT
        elif a_is_mom or b_is_mom:
            effective_tol = TOLERANCE_MOM
        else:
            effective_tol = TOLERANCE_OTHER
        return abs(fa - fb) <= effective_tol + 1e-9
    return False


def _suspected_causes(expected: str, actual: str) -> list[str]:
    """expected vs actual 차이로 의심 원인 추정 (공식 어디가 틀렸나)."""
    a = (actual or "").strip()
    if a in ("", "-"):
        return ["wrong filter (빈 결과)", "wrong dataset/column"]
    fe, fa = _to_float(expected), _to_float(actual)
    if fe is None or fa is None:
        return ["wrong cell/format"]
    is_mom = "%p" in (expected or "") or "%p" in (actual or "")
    causes: list[str] = []
    # 부호 반대 → period(mom) 의심
    if is_mom and (fe > 0) != (fa > 0):
        causes.append("wrong period")
    # 배수 관계 → scale/denominator 의심
    if fa != 0:
        r = fe / fa
        if 1.8 <= abs(r) <= 2.2 or 0.45 <= abs(r) <= 0.55:
            causes.append("wrong scale")
            causes.append("wrong denominator")
    # 기본: 분자/필터
    causes.append("wrong numerator")
    causes.append("wrong filter")
    # 중복 제거(순서 유지)
    seen, out = set(), []
    for c in causes:
        if c not in seen:
            seen.add(c); out.append(c)
    return out


def _infer_root_cause(ans_val: str, out_val: str) -> str:
    if not out_val or out_val.strip() in ("", "-"):
        return "missing_value"
    fa, fb = _to_float(ans_val), _to_float(out_val)
    if fa is not None and fb is not None:
        return "format_error" if abs(fa - fb) <= 0.15 else "wrong_value"
    return "wrong_cell"


def _gen_feedback(issues: list[VerificationIssue], route_to: str) -> str:
    """불일치 목록을 LLM에게 넘겨 구체적 피드백 생성."""
    api_key = get_anthropic_api_key()
    llm = ChatAnthropic(model="claude-sonnet-4-6", api_key=api_key, max_tokens=4000,
                        timeout=90, max_retries=2)

    issues_text = "\n".join(
        f"  slide={i.slide_idx} shape={i.shape_num} [{i.row},{i.col}] "
        f"key={i.value_key!r} expected={i.expected!r} actual={i.actual!r} "
        f"cause={i.root_cause}"
        for i in issues[:15]
    )
    if len(issues) > 15:
        issues_text += f"\n  ... 외 {len(issues) - 15}개"

    agent = "Calculator (계산식/필터 수정)" if route_to == "calculator" else "Filler (셀 위치/포맷 수정)"
    try:
        resp = llm.invoke([
            SystemMessage(content=_CONTRACT),
            HumanMessage(content=f"## 불일치 목록\n{issues_text}\n\n## 라우팅: {agent}\n\n구체적 수정 지시를 3-5줄로 작성하세요."),
        ])
        return resp.content.strip()
    except Exception:
        return issues_text


def verify_output(state: AgentState) -> dict:
    """SlideMapping.targets 기준 셀별 대조 → VerificationResult + 조건부 라우팅.

    반환: verification_result / validation_issues (override),
          retry_count / retry_feedback (override), messages (append)
    """
    mapping: SlideMapping | None = state.get("slide_mapping")
    output_path = state.get("execution_output_path") or state.get("output_path")
    answer_key_path = state.get("answer_key_path")
    retry_count = state.get("retry_count", 0)

    # ── 최대 재시도 선제 확인 (Manager가 주 제어, Verifier는 보조 확인) ─
    if retry_count >= MAX_RETRIES:
        result = VerificationResult(
            passed=False, route_to="end",
            feedback=f"최대 재시도 횟수({MAX_RETRIES}) 도달 — 종료",
        )
        print(f"[Verifier] {result.feedback}")
        return {
            "verification_result": result,
            "retry_count": retry_count,
            "pending_gate": "after_verify",
            "messages": [AIMessage(content=result.feedback, name="Verifier")],
        }

    # ── 출력 PPT 없음 ─────────────────────────────────────────────────
    if not output_path or not os.path.exists(output_path):
        result = VerificationResult(
            passed=False, route_to="filler", feedback="출력 PPT 없음 — Filler 재실행",
        )
        print("[Verifier] 출력 PPT 없음")
        return {
            "errors": ["출력 PPT 없음"],
            "verification_result": result,
            "pending_gate": "after_verify",
            "messages": [AIMessage(content="출력 PPT 없음", name="Verifier")],
        }

    # ── SlideMapping 없음 → 검증 불가 ────────────────────────────────
    if not mapping or not mapping.targets:
        result = VerificationResult(passed=True, route_to="end", feedback="SlideMapping 없음 — 검증 건너뜀")
        return {
            "verification_result": result,
            "pending_gate": "after_verify",
            "messages": [AIMessage(content=result.feedback, name="Verifier")],
        }

    # ── 정답지 없음 → 검증 건너뜀 ────────────────────────────────────
    if not answer_key_path or not os.path.exists(answer_key_path):
        result = VerificationResult(passed=True, route_to="end", feedback="정답지 없음 — 검증 건너뜀")
        return {
            "verification_result": result,
            "pending_gate": "after_verify",
            "messages": [AIMessage(content=result.feedback, name="Verifier")],
        }

    # ── Pipeline Skill 0: KPI Conservation Report ────────────────────────
    kpi_counts_val = state.get("kpi_counts") or {}
    if len(kpi_counts_val) >= 2:
        print_conservation_report(kpi_counts_val)
        cons_ok, cons_errors = check_conservation(kpi_counts_val)
        if not cons_ok:
            result = VerificationResult(
                passed=False, route_to="end",
                feedback="PIPELINE FAIL: KPI Conservation 위반 — " + " | ".join(cons_errors),
            )
            print(f"[Verifier] {result.feedback}")
            return {
                "verification_result": result,
                "errors": cons_errors,
                "messages": [AIMessage(content=result.feedback, name="Verifier")],
            }

    print(f"[Verifier] SlideMapping {len(mapping.targets)}개 셀 검증 중...")

    # ── 차트 완성도 검사: 출력 차트 중 모든 계열이 0/None이면 미채움 실패 ──────
    try:
        _out = Presentation(output_path)
        zero_charts = []
        for zi, zsl in enumerate(_out.slides):
            for zsh in zsl.shapes:
                if hasattr(zsh, "has_chart") and zsh.has_chart:
                    vv = [v for s in zsh.chart.series for v in s.values]
                    if vv and all(v in (0, None) for v in vv):
                        zero_charts.append(f"slide={zi} {zsh.name}")
        if zero_charts:
            print(f"[Verifier] ⚠ 미채움 차트 {len(zero_charts)}개 (모든 계열값 0): {zero_charts[:5]}")
    except Exception:
        zero_charts = []

    try:
        ans_prs = Presentation(answer_key_path)
        out_prs = Presentation(output_path)
    except Exception as e:
        new_count = retry_count + 1
        result = VerificationResult(
            passed=False,
            route_to="filler" if new_count < MAX_RETRIES else "end",
            feedback=f"PPT 로드 오류: {e}",
        )
        return {
            "errors": [f"검증 PPT 로드 실패: {e}"],
            "verification_result": result,
            "retry_count": new_count,
        }

    # ── SlideMapping.targets 기준 셀 대조 ─────────────────────────────
    issues: list[VerificationIssue] = []
    failure_reports: dict[str, dict] = {}   # key → FailureReport (key당 1개)
    ok_count = 0

    for t in mapping.targets:
        s_type = getattr(t, "shape_type", "table")
        if s_type == "text":
            ans_val = read_text_shape(ans_prs, t.slide_idx, t.shape_num,
                                      shape_id=t.shape_id) or ""
            out_val = read_text_shape(out_prs, t.slide_idx, t.shape_num,
                                      shape_id=t.shape_id) or ""
        else:
            ans_val = read_cell(ans_prs, t.slide_idx, t.shape_num, t.row, t.col,
                                shape_id=t.shape_id) or ""
            out_val = read_cell(out_prs, t.slide_idx, t.shape_num, t.row, t.col,
                                shape_id=t.shape_id) or ""

        if _is_close(ans_val, out_val):
            ok_count += 1
        else:
            cause = _infer_root_cause(ans_val, out_val)
            if cause == "wrong_value":
                suggestion = f"key={t.value_key!r}: 계산식/필터 재확인 (expected={ans_val!r} actual={out_val!r})"
            elif cause == "format_error":
                suggestion = f"key={t.value_key!r}: 포맷 수정 (expected={ans_val!r} actual={out_val!r})"
            elif cause == "missing_value":
                suggestion = f"slide={t.slide_idx} shape={t.shape_num} [{t.row},{t.col}] 값 누락"
            else:
                suggestion = f"slide={t.slide_idx} shape={t.shape_num} [{t.row},{t.col}] 셀 위치 확인"

            issues.append(VerificationIssue(
                slide_idx=t.slide_idx, shape_num=t.shape_num,
                row=t.row, col=t.col, value_key=t.value_key,
                expected=ans_val, actual=out_val,
                root_cause=cause, suggestion=suggestion,
            ))
            # 구조화 FailureReport (계산 오류 key만 — Planner 재합성 입력)
            if cause in ("wrong_value", "missing_value") and t.value_key and t.value_key not in failure_reports:
                failure_reports[t.value_key] = FailureReport(
                    key=t.value_key, expected=ans_val, actual=out_val,
                    suspected_causes=_suspected_causes(ans_val, out_val),
                    next_action="regenerate_formula_candidates",
                ).model_dump()

    total = len(mapping.targets)
    # 미채움 차트(전부 0)도 실패 조건에 포함
    passed = len(issues) == 0 and not zero_charts
    print(f"[Verifier] {ok_count}/{total} 일치, {len(issues)}개 불일치"
          + (f", 미채움 차트 {len(zero_charts)}개" if zero_charts else ""))

    # ── 통과 (셀 전부 일치 + 미채움 차트 없음) ──────────────────────
    if passed:
        result = VerificationResult(
            passed=True, total_checked=total, ok_count=ok_count,
            issues=[], route_to="end", feedback="검증 통과 ✓",
        )
        return {
            "verification_result": result,
            "pending_gate": "after_verify",
            "messages": [AIMessage(content=f"검증 통과: {ok_count}/{total}", name="Verifier")],
        }

    # ── 실패: 원인 분류 → route_to 결정 (실제 라우팅은 Manager가 담당) ───
    match_rate = ok_count / total if total > 0 else 0
    if match_rate < 0.90:
        # 90% 미만은 데이터 분포 차이 가능성 높음 → Manager에게 end 권고
        result = VerificationResult(
            passed=False, total_checked=total, ok_count=ok_count,
            issues=issues, route_to="end",
            feedback=f"일치율 {match_rate:.0%} — 데이터 분포 차이로 추정",
        )
        msg = f"검증 완료: {ok_count}/{total} ({match_rate:.0%}) | 90% 미달 → Manager에 end 권고"
        print(f"[Verifier] {msg}")
        return {
            "verification_result": result,
            "validation_issues": [
                ValidationIssue(slide_idx=i.slide_idx, shape_num=i.shape_num,
                                 row=i.row, col=i.col, expected=i.expected, actual=i.actual)
                for i in issues
            ],
            "pending_gate": "after_verify",
            "messages": [AIMessage(content=msg, name="Verifier")],
        }

    # wrong_value(값 오류) + missing_value(빈 값)는 명세/계산 문제 → "calculator"
    #   → Manager가 이를 Planner 자아 진화로 라우팅한다.
    # wrong_cell·format_error(셀 위치/포맷)는 "filler" 소관.
    spec_issues = [i for i in issues if i.root_cause in ("wrong_value", "missing_value")]
    route_to = "calculator" if len(spec_issues) >= len(issues) * 0.5 else "filler"

    print(f"[Verifier] LLM 피드백 생성 → route_to={route_to} (Manager가 최종 결정)")
    feedback = _gen_feedback(issues, route_to)

    result = VerificationResult(
        passed=False, total_checked=total, ok_count=ok_count,
        issues=issues, route_to=route_to, feedback=feedback,
    )
    msg = f"검증 실패: {ok_count}/{total} | {len(issues)}개 불일치 → {route_to} 권고"
    print(f"[Verifier] {msg}")

    # 공유 메모리: validation_events 기록 (실패 이력 누적)
    try:
        from core.predefined.shared_memory import get_memory
        get_memory().log_validations_bulk([
            (i.value_key, i.expected, i.actual, i.root_cause, route_to) for i in issues
        ])
        get_memory().log_message("Verifier", "Planner",
                                 f"검증 {ok_count}/{total}, FailureReport {len(failure_reports)}건",
                                 ["failure_reports"])
        print(f"[Verifier] 공유 메모리: validation_events {len(issues)}건 기록")
    except Exception as e:
        print(f"[Verifier] 메모리 기록 실패(무시): {e}")

    print(f"[Verifier] FailureReport {len(failure_reports)}건 생성 → Planner 재합성 입력")
    return {
        "verification_result": result,
        "validation_issues": [
            ValidationIssue(slide_idx=i.slide_idx, shape_num=i.shape_num,
                             row=i.row, col=i.col, expected=i.expected, actual=i.actual)
            for i in issues
        ],
        "failure_reports": list(failure_reports.values()),
        "retry_feedback": feedback,
        "pending_gate": "after_verify",
        "messages": [AIMessage(content=msg, name="Verifier")],
    }
