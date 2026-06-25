"""Insight Writer 에이전트 — 채워진 KPI 값을 근거로 슬라이드 인사이트 문장 자동 작성.

Filler가 KPI를 모두 채운 뒤(execution_output_path) 실행한다. 템플릿의 인사이트
placeholder({{...}})를 찾아, 그 슬라이드에 실제로 채워진 수치를 요약한 문장으로 교체한다.

[입력] execution_output_path(채워진 PPT), slide_mapping, calculation_result
[출력] 같은 PPT에 인사이트 문장 기록 (execution_output_path 그대로 유지)
대상/규칙: insight_targets.json(같은 폴더) + insight_writer-skills.md
"""
from __future__ import annotations

import json
import os
import re
from collections import defaultdict

from langchain_anthropic import ChatAnthropic
from langchain_core.messages import AIMessage, HumanMessage, SystemMessage
from pptx import Presentation

from agents.state import AgentState
from agents.models import CalculationResult, SlideMapping
from agents.utils import get_anthropic_api_key, load_contract, load_skills
from core.predefined.pptx_scanner import scan_pptx, write_insight_shape

_DIR = os.path.dirname(os.path.abspath(__file__))
_CONFIG_PATH = os.path.join(_DIR, "insight_targets.json")

_CONTRACT = load_contract("insight_writer")
_SKILLS = load_skills("insight_writer")

_llm: ChatAnthropic | None = None


def _get_llm() -> ChatAnthropic:
    global _llm
    if _llm is None:
        _llm = ChatAnthropic(
            model="claude-sonnet-4-6",
            api_key=get_anthropic_api_key(),
            max_tokens=1024,
            timeout=60,       # 문장 1개 생성 — 짧음
            max_retries=2,
        )
    return _llm


def _load_config() -> dict:
    try:
        with open(_CONFIG_PATH, encoding="utf-8") as f:
            return json.load(f)
    except Exception:
        return {"enabled": True, "placeholder_patterns": [r"\{\{[^}]+\}\}"]}


def _iter_text_shapes(shapes: list[dict]):
    """스캔 결과에서 텍스트 도형을 (그룹 재귀) 평탄화."""
    for sh in shapes:
        if sh.get("type") == "group":
            yield from _iter_text_shapes(sh.get("children", []))
        elif sh.get("type") == "text" and sh.get("text"):
            yield sh


def _slide_kpi_context(mapping: SlideMapping,
                       calc: CalculationResult) -> dict[int, list[str]]:
    """slide_idx → ["지표설명: 포맷값", ...] (채워진 숫자 KPI만; text 자리 제외)."""
    out: dict[int, list[str]] = defaultdict(list)
    for t in mapping.targets:
        if t.format_type == "text":
            continue
        cv = calc.get(t.value_key)
        if cv and cv.raw_value is not None:
            label = t.context or t.value_key
            out[t.slide_idx].append(f"- {label}: {cv.formatted_value}")
    return out


def _gen_insight(slide_idx: int, placeholder: str, kpi_lines: list[str]) -> str:
    """슬라이드 KPI 값으로 인사이트 문장 1개 생성 (LLM)."""
    system = _CONTRACT + (f"\n\n---\n## Insight Writer Skills\n{_SKILLS}" if _SKILLS else "")
    human = (
        f"## 슬라이드 {slide_idx} 인사이트 자리\n"
        f"placeholder: {placeholder!r}\n\n"
        f"## 이 슬라이드에 채워진 KPI 값 (오직 이 수치만 사용)\n"
        + "\n".join(kpi_lines)
        + "\n\n위 수치만 근거로 이 슬라이드를 요약하는 한국어 인사이트 문장을 작성하라.\n"
        "- placeholder가 짧으면 1문장, 요약 영역이면 2~3문장\n"
        "- MoM/vs 비교로 방향성(상승/하락/우위/열세) 명시, 수치 인용\n"
        "- 없는 수치 지어내기 금지, 다른 슬라이드 내용 혼합 금지\n"
        "- 문장만 출력 (레이블·따옴표·설명·마크다운 없이)"
    )
    resp = _get_llm().invoke([
        SystemMessage(content=system),
        HumanMessage(content=human),
    ])
    return (resp.content or "").strip().strip('"').strip()


def write_insights(state: AgentState) -> dict:
    """채워진 PPT의 인사이트 placeholder를 슬라이드 수치 요약 문장으로 교체."""
    out_path = state.get("execution_output_path") or state.get("output_path")
    mapping: SlideMapping | None = state.get("slide_mapping")
    calc: CalculationResult | None = state.get("calculation_result")

    def _done(msg: str) -> dict:
        print(f"[InsightWriter] {msg}")
        return {
            "pending_gate": "after_insight",
            "messages": [AIMessage(content=msg, name="InsightWriter")],
        }

    cfg = _load_config()
    if not cfg.get("enabled", True):
        return _done("비활성화됨 — 건너뜀")
    if not out_path or not os.path.exists(out_path):
        return _done("채워진 PPT 없음 — 건너뜀")
    if not mapping or not calc:
        return _done("매핑/계산결과 없음 — 건너뜀")
    if not get_anthropic_api_key():
        return _done("API 키 없음 — 인사이트 생략")

    patterns = [re.compile(p) for p in cfg.get("placeholder_patterns", [r"\{\{[^}]+\}\}"])]

    # 1. 채워진 PPT 스캔 → 인사이트 대상 탐색
    #    A) {{...}} 패턴 placeholder
    #    B) "Insight & Summary" 레이블 옆 빈 텍스트 박스
    prs = Presentation(out_path)
    targets: list[tuple[int, int | None, int | None, str]] = []

    # A) {{...}} 패턴
    scan = scan_pptx(out_path, read_values=True)
    for sl in scan:
        for sh in _iter_text_shapes(sl["shapes"]):
            text = sh["text"]
            if any(p.search(text) for p in patterns):
                targets.append((sl["slide_idx"], sh.get("num"), sh.get("shape_id"), text))

    # B) "Insight & Summary" 레이블 인접 빈 텍스트 박스
    #    레이블 위치(top) 기준 ±1.5cm(540000 EMU) 이내 + 가장 넓은 빈 TextBox
    _INSIGHT_LABEL = re.compile(r"insight\s*[&＆]\s*summary", re.I)
    _INSIGHT_TOP_TOL = 540_000  # EMU (~1.5cm)
    for sidx, slide in enumerate(prs.slides):
        # 이 슬라이드에 이미 {{}} 대상이 있으면 스킵 (중복 방지)
        if any(t[0] == sidx for t in targets):
            continue
        label_shape = None
        for sh in slide.shapes:
            if hasattr(sh, "has_text_frame") and sh.has_text_frame:
                if _INSIGHT_LABEL.search(sh.text_frame.text or ""):
                    label_shape = sh
                    break
        if not label_shape:
            continue
        # label 근처 빈 TextBox 중 가장 넓은 것
        best = None
        best_width = 0
        for sh in slide.shapes:
            if not (hasattr(sh, "has_text_frame") and sh.has_text_frame):
                continue
            if (sh.text_frame.text or "").strip():
                continue  # 비어있어야 함
            if sh.shape_id == label_shape.shape_id:
                continue
            top_diff = abs((sh.top or 0) - (label_shape.top or 0))
            if top_diff > _INSIGHT_TOP_TOL:
                continue
            if (sh.width or 0) > best_width:
                best_width = sh.width or 0
                best = sh
        if best:
            targets.append((sidx, None, best.shape_id,
                            f"[Insight & Summary — slide {sidx}]"))
            print(f"[InsightWriter] 'Insight & Summary' 빈 박스 발견: "
                  f"slide={sidx} shape_id={best.shape_id}")

    if not targets:
        return _done("인사이트 자리 없음 — 작성 대상 0건")

    # 2. 슬라이드별 채워진 KPI 컨텍스트
    slide_kpis = _slide_kpi_context(mapping, calc)

    # 3. 인사이트 생성 + 기록 (근거 수치 없으면 생략 — 환각 방지)
    print(f"[InsightWriter] 인사이트 자리 {len(targets)}건 발견 — 작성 시작")
    written = skipped = 0
    for slide_idx, shape_num, shape_id, ph in targets:
        # 이 슬라이드 + 인접 슬라이드의 KPI 합산 (인사이트 슬라이드가 summary인 경우)
        kpis = slide_kpis.get(slide_idx, [])
        if not kpis:
            skipped += 1
            continue
        try:
            sentence = _gen_insight(slide_idx, ph, kpis)
        except Exception as e:
            print(f"[InsightWriter]   slide {slide_idx}: LLM 실패({str(e)[:60]}) — 생략")
            skipped += 1
            continue
        if sentence and write_insight_shape(prs, slide_idx, shape_num, sentence,
                                           shape_id=shape_id):
            written += 1
            print(f"[InsightWriter]   slide {slide_idx} ← {sentence[:60]}…")
        else:
            skipped += 1

    if written:
        prs.save(out_path)

    return _done(f"인사이트 작성 완료: {written}건 기록, {skipped}건 생략 (대상 {len(targets)})")
