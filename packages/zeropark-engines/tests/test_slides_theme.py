"""Slide theme rendering: themed pptx opens and resolve_theme merges overrides."""

import io

import pytest
from pptx import Presentation

from zeropark_core.capabilities import Capability
from zeropark_core.models import TaskRequest
from zeropark_core.store import LocalArtifactStore
from zeropark_engines.slides import PptxSlidesEngine, THEMES, resolve_theme


def test_resolve_theme_named_and_overrides():
    assert resolve_theme({}) == THEMES["default"]
    assert resolve_theme({"theme": "dark"})["background"] == "#111827"
    merged = resolve_theme({"theme": "corporate", "theme_overrides": {"accent": "#FF0000", "bogus": "x"}})
    assert merged["accent"] == "#FF0000"
    assert "bogus" not in merged


@pytest.mark.asyncio
async def test_themed_deck_renders_valid_pptx(tmp_path):
    store = LocalArtifactStore(base_dir=str(tmp_path))
    engine = PptxSlidesEngine(store=store)
    task = TaskRequest(
        prompt="deck",
        capability=Capability.SLIDES,
        params={
            "title": "Themed Deck",
            "theme": "dark",
            "outline": [
                {"title": "Intro", "bullets": ["one", "two"], "notes": "say hello"},
                {"title": "Part 1", "layout": "section"},
            ],
        },
    )
    result = await engine.cap_slides(task, "t1")
    deck_path = result.artifacts[0].uri
    prs = Presentation(deck_path)
    assert len(prs.slides.__iter__().__self__._sldIdLst) == 3  # title + 2 content
    # speaker notes survived
    notes = prs.slides[1].notes_slide.notes_text_frame.text
    assert notes == "say hello"
