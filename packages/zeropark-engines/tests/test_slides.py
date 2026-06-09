from __future__ import annotations

import asyncio

from pptx import Presentation
from zeropark_core.capabilities import Capability
from zeropark_core.models import TaskRequest
from zeropark_core.store import LocalArtifactStore

from zeropark_engines.slides import PptxSlidesEngine


def test_slides_renders_real_pptx(tmp_path) -> None:
    store = LocalArtifactStore(base_dir=str(tmp_path))
    engine = PptxSlidesEngine(store=store)
    task = TaskRequest(
        prompt="Quarterly review",
        capability=Capability.SLIDES,
        params={
            "title": "Q2 Review",
            "outline": [
                {"title": "Revenue", "bullets": ["Up 12%", "New region"]},
                {"title": "Risks", "bullets": ["Hiring"]},
            ],
        },
    )
    result = asyncio.run(engine.execute(task, task_id="deck1"))
    assert result.status.value == "succeeded"
    deck = result.artifacts[0]
    assert deck.kind == "deck"

    saved = tmp_path / "deck1.pptx"
    assert saved.exists()
    prs = Presentation(str(saved))
    # title slide + 2 content slides
    assert len(prs.slides) == 3
