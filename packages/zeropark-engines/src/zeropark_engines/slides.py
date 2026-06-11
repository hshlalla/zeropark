"""SLIDES — render an outline to a real .pptx. Native (python-pptx).

Design reference: Presenton (Apache-2.0). No Presenton code or service is used.
Content generation (turning a prompt into an outline) is a pluggable LLM step;
this engine renders whatever outline it is given, so it works with zero config.
"""

from __future__ import annotations

import asyncio
import json
import re
import io
import os
import shutil
import subprocess
import tempfile
from collections.abc import AsyncIterator
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt
from zeropark_core.capabilities import Capability
from zeropark_core import ArtifactStore
from zeropark_core.models import Artifact, RunEvent, TaskRequest, TaskResult, TaskStatus
from zeropark_core.llm import BaseLLMClient, ChatMessage

from zeropark_engines.base import NativeEngine

_PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


def _hex_to_rgb(value: str) -> RGBColor:
    value = value.lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


# Built-in themes (Presenton-style template layer). A deployment can also pass
# a fully custom theme dict in task params — e.g. matched to client branding.
THEMES: dict[str, dict[str, Any]] = {
    "default": {
        "background": "#FFFFFF",
        "title_color": "#1F2937",
        "body_color": "#374151",
        "accent": "#4F46E5",
        "font": "Calibri",
    },
    "dark": {
        "background": "#111827",
        "title_color": "#F9FAFB",
        "body_color": "#D1D5DB",
        "accent": "#818CF8",
        "font": "Calibri",
    },
    "corporate": {
        "background": "#F8FAFC",
        "title_color": "#0F172A",
        "body_color": "#334155",
        "accent": "#0EA5E9",
        "font": "Arial",
    },
}


def resolve_theme(params: dict[str, Any]) -> dict[str, Any]:
    """Resolve a theme: named theme + optional per-key overrides (e.g. accent
    set to the deployment's branding primary color)."""
    theme = dict(THEMES.get(str(params.get("theme", "default")), THEMES["default"]))
    overrides = params.get("theme_overrides") or {}
    theme.update({k: v for k, v in overrides.items() if k in theme})
    return theme

def _download_image(url: str) -> io.BytesIO | None:
    try:
        if url.startswith("http"):
            import httpx
            resp = httpx.get(url, timeout=10)
            resp.raise_for_status()
            return io.BytesIO(resp.content)
        elif os.path.exists(url):
            with open(url, "rb") as f:
                return io.BytesIO(f.read())
    except Exception:
        pass
    return None


class PptxSlidesEngine(NativeEngine):
    id = "zeropark_engines.slides.pptx"
    name = "Pptx Renderer Engine"
    capabilities = {Capability.SLIDES}
    reference = "Presenton (Apache-2.0) - design reference only"

    def __init__(self, store: ArtifactStore) -> None:
        self.store = store

    @staticmethod
    def _paint_background(slide, theme: dict) -> None:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = _hex_to_rgb(theme["background"])

    @staticmethod
    def _style_text_frame(text_frame, *, color: str, font: str, size: int, bold: bool = False) -> None:
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = _hex_to_rgb(color)
                run.font.name = font
                run.font.size = Pt(size)
                run.font.bold = bold

    def _render_pptx(
        self, title: str, subtitle: str, outline: list[dict], theme: dict | None = None, theme_name: str = "default"
    ) -> bytes:
        theme = theme or THEMES["default"]
        
        template_dir = Path(__file__).parent / "templates"
        template_path = template_dir / f"{theme_name}.pptx"
        if template_path.exists():
            prs = Presentation(str(template_path))
        else:
            prs = Presentation()

        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        self._paint_background(title_slide, theme)
        if title_slide.shapes.title:
            title_slide.shapes.title.text = title
            self._style_text_frame(
                title_slide.shapes.title.text_frame,
                color=theme["accent"], font=theme["font"], size=40, bold=True,
            )
        if len(title_slide.placeholders) > 1:
            title_slide.placeholders[1].text = subtitle
            self._style_text_frame(
                title_slide.placeholders[1].text_frame,
                color=theme["body_color"], font=theme["font"], size=18,
            )

        for item in outline:
            image_url = item.get("image")
            if image_url:
                layout_index = 8 if len(prs.slide_layouts) > 8 else 1
            elif item.get("layout") == "section":
                layout_index = 2
            else:
                layout_index = 1
                
            slide = prs.slides.add_slide(prs.slide_layouts[layout_index])
            self._paint_background(slide, theme)
            
            if slide.shapes.title:
                slide.shapes.title.text = item.get("title", "")
                self._style_text_frame(
                    slide.shapes.title.text_frame,
                    color=theme["title_color"], font=theme["font"], size=30, bold=True,
                )
            
            bullets = item.get("bullets", []) or []
            
            body_placeholder = None
            picture_placeholder = None
            for shape in slide.placeholders:
                if shape.placeholder_format.type == 1:
                    continue
                if shape.placeholder_format.type == 18:
                    picture_placeholder = shape
                elif body_placeholder is None:
                    body_placeholder = shape

            if bullets and body_placeholder:
                body = body_placeholder.text_frame
                body.text = str(bullets[0])
                for bullet in bullets[1:]:
                    body.add_paragraph().text = str(bullet)
                self._style_text_frame(
                    body, color=theme["body_color"], font=theme["font"], size=18,
                )

            if image_url:
                img_io = _download_image(image_url)
                if img_io:
                    try:
                        if picture_placeholder:
                            picture_placeholder.insert_picture(img_io)
                        else:
                            from pptx.util import Inches
                            slide.shapes.add_picture(img_io, left=Inches(5), top=Inches(2), width=Inches(4))
                    except Exception:
                        pass

            notes = item.get("notes")
            if notes:
                # accessing notes_slide creates it when absent
                slide.notes_slide.notes_text_frame.text = str(notes)

        bio = io.BytesIO()
        prs.save(bio)
        return bio.getvalue()

    async def cap_slides(self, task: TaskRequest, task_id: str) -> TaskResult:
        outline = task.params.get("outline") or [
            {"title": task.params.get("title") or task.prompt[:80], "bullets": []}
        ]
        deck_title = task.params.get("title") or outline[0].get("title") or task.prompt[:80]
        subtitle = task.params.get("subtitle", "Generated by Zeropark")

        theme_name = str(task.params.get("theme", "default"))
        theme = resolve_theme(task.params)
        pptx_bytes = self._render_pptx(deck_title, subtitle, outline, theme=theme, theme_name=theme_name)
        
        # Save to store
        filename = f"{task_id}.pptx"
        file_uri = self.store.save(filename, pptx_bytes)

        n_slides = len(outline) + 1
        artifacts = [Artifact(
            id=f"{task_id}_deck",
            kind="deck",
            title=deck_title,
            mime_type=_PPTX_MIME,
            uri=file_uri,
            metadata={"slides": n_slides},
        )]
        
        # Convert to PDF
        soffice_path = shutil.which("soffice") or shutil.which("libreoffice")
        if soffice_path:
            try:
                temp_dir = Path(tempfile.gettempdir())
                temp_pptx = temp_dir / f"{task_id}.pptx"
                temp_pptx.write_bytes(pptx_bytes)
                
                # Execute soffice headlessly
                subprocess.run(
                    [soffice_path, "--headless", "--convert-to", "pdf", str(temp_pptx), "--outdir", str(temp_dir)],
                    check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL
                )
                temp_pdf = temp_dir / f"{task_id}.pdf"
                
                if temp_pdf.exists():
                    pdf_bytes = temp_pdf.read_bytes()
                    pdf_uri = self.store.save(f"{task_id}.pdf", pdf_bytes)
                    artifacts.append(Artifact(
                        id=f"{task_id}_pdf",
                        kind="file",
                        title=f"{deck_title} (PDF)",
                        mime_type="application/pdf",
                        uri=pdf_uri,
                        metadata={"slides": n_slides},
                    ))
                    temp_pdf.unlink(missing_ok=True)
                temp_pptx.unlink(missing_ok=True)
            except Exception as e:
                print(f"PDF conversion failed: {e}")

        return TaskResult(
            task_id=task_id,
            status=TaskStatus.SUCCEEDED,
            capability=Capability.SLIDES,
            provider_id=self.id,
            artifacts=artifacts,
            metrics={"slides": n_slides},
        )

class LLMSlidesEngine(NativeEngine):
    id = "llm-slides"
    name = "LLM Slides Generator"
    capabilities = frozenset({Capability.SLIDES})
    reference = "Presenton (Apache-2.0) - LLM Outline Generation"

    def __init__(self, llm_client: BaseLLMClient, renderer: PptxSlidesEngine, model_name: str = "gpt-4o-mini") -> None:
        self.llm_client = llm_client
        self.renderer = renderer
        self.model_name = model_name

    async def _generate_outline(self, prompt: str) -> tuple[str, list[dict]]:
        system_msg = ChatMessage(
            role="system",
            content=(
                "You are an expert presentation designer. Create an outline for a slide deck based on the user's prompt. "
                "Also, intelligently select a suitable design theme based on the context of the prompt. "
                "Available themes: ['default', 'dark', 'corporate']. (e.g. 'dark' for tech/night/hacking, 'corporate' for business/finance, 'default' for general). "
                "Your output MUST be a strict JSON object with this exact structure: "
                "{\"theme\": \"corporate\", \"slides\": [{\"title\": \"Slide Title\", \"bullets\": [\"Point 1\"], \"notes\": \"speaker notes\"}]}. "
                "Generate 4 to 8 slides with concise, punchy bullets (max 8 words each). "
                "Return ONLY the JSON object, nothing else."
            )
        )
        user_msg = ChatMessage(role="user", content=prompt)
        
        loop = asyncio.get_event_loop()
        response = await loop.run_in_executor(
            None,
            lambda: self.llm_client.chat_completion([system_msg, user_msg], self.model_name, temperature=0.5)
        )
        
        text = response.content.strip()
        # Find JSON object using regex if there's markdown wrapping
        match = re.search(r'\{.*\}', text, re.DOTALL)
        if match:
            text = match.group(0)
            
        try:
            data = json.loads(text)
            if isinstance(data, dict) and "slides" in data:
                return data.get("theme", "default"), data["slides"]
        except Exception:
            pass
            
        # Fallback if parsing fails
        return "default", [{"title": prompt[:80], "bullets": ["Could not generate detailed outline."]}]

    async def cap_slides(self, task: TaskRequest, task_id: str) -> TaskResult:
        # 1. Generate Theme & Outline via LLM
        theme, outline = await self._generate_outline(task.prompt)
        
        # 2. Inject outline and theme into task params
        if not task.params:
            task.params = {}
        task.params["outline"] = outline
        task.params["theme"] = theme
        
        # 3. Delegate to renderer
        return await self.renderer.cap_slides(task, task_id)

    async def stream(self, task: TaskRequest, *, task_id: str) -> AsyncIterator[RunEvent]:
        yield RunEvent(
            type="status",
            task_id=task_id,
            provider_id=self.id,
            message="started",
            data={"capability": task.capability.value},
        )
        yield RunEvent(
            type="log",
            task_id=task_id,
            provider_id=self.id,
            message="[Planner] Designing slide outline and selecting theme based on your request..."
        )
        
        try:
            # 1. Generate Theme & Outline via LLM
            theme, outline = await self._generate_outline(task.prompt)
            yield RunEvent(
                type="log",
                task_id=task_id,
                provider_id=self.id,
                message=f"[Planner] Generated {len(outline)} slides with theme '{theme}'."
            )
            
            # 2. Inject outline and theme into task params
            if not task.params:
                task.params = {}
            task.params["outline"] = outline
            task.params["theme"] = theme
            
            yield RunEvent(
                type="log",
                task_id=task_id,
                provider_id=self.id,
                message="[Renderer] Generating PowerPoint (.pptx) and PDF files..."
            )
            
            # 3. Delegate to renderer
            result = await self.renderer.cap_slides(task, task_id)
            
            for artifact in result.artifacts:
                yield RunEvent(
                    type="artifact",
                    task_id=task_id,
                    provider_id=self.id,
                    data={"artifact": artifact.model_dump(mode="json")},
                )
            yield RunEvent(
                type="done",
                task_id=task_id,
                provider_id=self.id,
                data={"status": result.status.value, "result": result.model_dump(mode="json")},
            )
        except Exception as exc:
            yield RunEvent(
                type="error",
                task_id=task_id,
                provider_id=self.id,
                message=str(exc)
            )
